"""
使用AI增强的PPT转视频服务
"""
from services.ppt_to_video import PPTToVideoService
from services.ai_service import TTSService, DigitalHumanService
from pptx import Presentation
import asyncio
import os


class AIPTTToVideoService(PPTToVideoService):
    """AI增强的PPT转视频服务"""
    
    def __init__(self, tts_provider: str = "baidu", avatar_provider: str = "tencent"):
        super().__init__()
        self.ai_enabled = False
        self.tts_service = None
        try:
            self.tts_service = TTSService(provider=tts_provider)
            # 数字人服务已禁用，不需要初始化
            # self.avatar_service = DigitalHumanService(provider=avatar_provider)
            self.ai_enabled = True
            print("[INFO] TTS服务初始化成功")
        except Exception as e:
            print(f"[WARN] TTS服务初始化失败，将使用基础功能（无音频）: {str(e)}")
            import traceback
            traceback.print_exc()
            self.ai_enabled = False
    
    async def convert_to_video(
        self,
        ppt_path: str,
        ai_avatar_enabled: bool = True
    ) -> str:
        """
        将PPT转换为视频（AI增强版）
        """
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None,
            self._convert_to_video_sync_ai,
            ppt_path,
            ai_avatar_enabled
        )
    
    def _convert_to_video_sync_ai(
        self,
        ppt_path: str,
        ai_avatar_enabled: bool
    ) -> str:
        """同步转换PPT为视频（AI增强）"""
        try:
            import traceback
            from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
            from moviepy.video.io.VideoFileClip import VideoFileClip
            from moviepy.video.compositing.CompositeVideoClip import CompositeVideoClip
            import imageio_ffmpeg
            
            print(f"开始转换PPT: {ppt_path}")
            
            # 加载PPT
            prs = Presentation(ppt_path)
            print(f"PPT加载成功，共 {len(prs.slides)} 张幻灯片")
            
            # 提取PPT文本内容
            ppt_text = self._extract_text_from_ppt(prs)
            print(f"提取PPT文本: {ppt_text[:100]}...")
            
            # 将每张幻灯片转换为图片
            slide_images = []
            for i, slide in enumerate(prs.slides):
                print(f"转换幻灯片 {i+1}/{len(prs.slides)}")
                img_path = self._slide_to_image(slide, i)
                slide_images.append(img_path)
            
            print(f"图片生成完成，共 {len(slide_images)} 张")
            
            # 创建视频片段
            print("开始创建视频片段...")
            try:
                # 验证图片文件是否存在
                for img_path in slide_images:
                    if not os.path.exists(img_path):
                        raise Exception(f"图片文件不存在: {img_path}")
                
                print(f"验证完成，共 {len(slide_images)} 张图片")
                final_video = ImageSequenceClip(
                    slide_images,
                    durations=[self.slide_duration] * len(slide_images)
                )
                print("ImageSequenceClip创建成功")
                # 调整视频尺寸（尝试多种方法）
                try:
                    # 方法1: 尝试使用 vfx.resize
                    from moviepy.video.fx.all import resize
                    final_video = final_video.fx(resize, height=720)
                    print("视频尺寸调整完成")
                except Exception as e1:
                    try:
                        # 方法2: 尝试直接使用 resize
                        final_video = final_video.resize(height=720)
                        print("视频尺寸调整完成（方法2）")
                    except Exception as e2:
                        print(f"[WARN] 调整视频尺寸失败，使用原始尺寸: {str(e1)} / {str(e2)}")
                        # 如果调整尺寸失败，继续使用原始尺寸
                print("视频片段创建成功")
            except Exception as clip_error:
                print(f"[ERROR] 创建视频片段失败: {str(clip_error)}")
                import traceback
                traceback.print_exc()
                # 尝试使用基础服务
                print("尝试使用基础服务...")
                from services.ppt_to_video import PPTToVideoService
                base_service = PPTToVideoService()
                return base_service._convert_to_video_sync(ppt_path, False)
            
            # 如果启用AI增强（仅TTS，不包含数字人）
            # 暂时禁用AI功能，先确保基础视频生成正常
            use_ai_audio = False  # 暂时禁用，避免AI服务导致失败
            if use_ai_audio and ai_avatar_enabled and self.ai_enabled and self.tts_service:
                print("开始生成AI语音讲解...")
                # 使用异步方式生成TTS
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
                
                try:
                    # 生成TTS音频
                    print("生成语音...")
                    if not ppt_text or len(ppt_text.strip()) == 0:
                    print("[WARN] PPT文本为空，跳过语音生成")
                    else:
                        audio_path = loop.run_until_complete(
                            self.tts_service.generate_speech(ppt_text[:500])  # 限制长度避免TTS失败
                        )
                        print(f"语音生成完成: {audio_path}")
                        
                        # 将音频添加到视频
                        try:
                            from moviepy.audio.io.AudioFileClip import AudioFileClip
                            
                            audio_clip = AudioFileClip(audio_path)
                            
                            # 确保音频长度匹配视频长度
                            if audio_clip.duration > final_video.duration:
                                audio_clip = audio_clip.subclip(0, final_video.duration)
                            elif audio_clip.duration < final_video.duration:
                                # 如果音频较短，循环播放
                                audio_clip = audio_clip.loop(duration=final_video.duration)
                            
                            # 添加音频到视频
                            final_video = final_video.set_audio(audio_clip)
                            print("音频添加完成")
                        except Exception as audio_error:
                            print(f"添加音频失败: {str(audio_error)}")
                            import traceback
                            traceback.print_exc()
                            # 继续使用无音频视频
                    
                except Exception as e:
                    print(f"AI语音功能失败，使用无音频视频: {str(e)}")
                    import traceback
                    traceback.print_exc()
                    # 如果AI功能失败，继续使用无音频视频
                finally:
                    try:
                        loop.close()
                    except:
                        pass
            else:
                if not self.ai_enabled:
                    print("AI服务未启用，生成无音频视频")
                else:
                    print("跳过AI语音生成")
            
            # 输出视频路径
            output_path = self.temp_dir / f"video_ai_{os.path.basename(ppt_path)}.mp4"
            print(f"开始导出视频到: {output_path}")
            
            # 设置FFmpeg路径
            try:
                ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
                os.environ['IMAGEIO_FFMPEG_EXE'] = ffmpeg_path
                print(f"FFmpeg路径: {ffmpeg_path}")
            except Exception as ffmpeg_error:
                print(f"[WARN] 获取FFmpeg路径失败: {str(ffmpeg_error)}")
                # 继续尝试，可能系统PATH中有FFmpeg
            
            # 导出视频
            try:
                print(f"准备导出视频，路径: {output_path}")
                print(f"视频时长: {final_video.duration}秒, 是否有音频: {final_video.audio is not None}")
                
                # 关闭视频资源前导出
                final_video.write_videofile(
                    str(output_path),
                    fps=24,
                    codec='libx264',
                    audio_codec='aac' if final_video.audio else None,
                    temp_audiofile=str(self.temp_dir / "temp_audio.m4a") if final_video.audio else None,
                    remove_temp=True,
                    ffmpeg_params=['-y'],
                    threads=4  # 使用多线程加速
                )
                print(f"[INFO] 视频导出成功: {output_path}")
            except Exception as write_error:
                            print(f"[ERROR] 视频写入失败: {str(write_error)}")
                import traceback
                traceback.print_exc()
                # 尝试使用基础方法（无音频）
                try:
                    print("尝试使用基础方法导出（无音频）...")
                    # 移除音频
                    if final_video.audio:
                        final_video = final_video.without_audio()
                    final_video.write_videofile(
                        str(output_path),
                        fps=24,
                        codec='libx264',
                        remove_temp=True,
                        ffmpeg_params=['-y']
                    )
                    print(f"[INFO] 使用基础方法导出成功")
                except Exception as fallback_error:
                    print(f"[ERROR] 基础方法也失败: {str(fallback_error)}")
                    raise Exception(f"视频导出失败: {str(write_error)}")
            finally:
                # 确保释放视频资源
                try:
                    if 'final_video' in locals():
                        final_video.close()
                except:
                    pass
            
            # 清理临时文件
            for img_path in slide_images:
                if os.path.exists(img_path):
                    try:
                        os.unlink(img_path)
                    except:
                        pass
            
            return str(output_path)
            
        except Exception as e:
            import traceback
            error_trace = traceback.format_exc()
            print(f"[ERROR] 转换视频失败: {str(e)}")
            print(f"错误堆栈: {error_trace}")
            
            # 清理资源
            try:
                if 'final_video' in locals() and final_video:
                    final_video.close()
            except:
                pass
            
            # 清理临时图片文件
            if 'slide_images' in locals():
                for img_path in slide_images:
                    if os.path.exists(img_path):
                        try:
                            os.unlink(img_path)
                        except:
                            pass
            
            # 重新抛出异常，让上层处理
            raise Exception(f"转换视频失败: {str(e)}")
    
    def _extract_text_from_ppt(self, prs: Presentation) -> str:
        """从PPT中提取所有文本内容（增强版，支持多种编码和完整提取）"""
        text_parts = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            
            # 提取标题
            if slide.shapes.title:
                title_text = self._extract_text_from_shape(slide.shapes.title)
                if title_text:
                    slide_texts.append(f"标题: {title_text}")
            
            # 提取所有形状的文本（包括占位符）
            for shape_idx, shape in enumerate(slide.shapes):
                if shape != slide.shapes.title:
                    shape_text = self._extract_text_from_shape(shape)
                    if shape_text:
                        slide_texts.append(shape_text)
            
            # 如果该幻灯片有文本，添加到总文本中
            if slide_texts:
                text_parts.append(f"\n--- 幻灯片 {slide_idx + 1} ---\n")
                text_parts.extend(slide_texts)
        
        # 合并文本并清理
        full_text = "\n".join(text_parts)
        # 清理文本：移除控制字符，保留可打印字符
        cleaned_text = self._clean_text(full_text)
        
        # 如果提取的文本为空或太短，尝试备用方法
        if not cleaned_text or len(cleaned_text.strip()) < 10:
            print("[WARN] 主要提取方法未获取到文本，尝试备用方法...")
            cleaned_text = self._extract_text_fallback(prs)
        
        return cleaned_text
    
    def _extract_text_from_shape(self, shape) -> str:
        """从形状中递归提取文本（增强版，支持多种编码和完整提取）"""
        text_parts = []
        
        try:
            # 方法1: 从文本框中提取段落（最完整的方法）
            if hasattr(shape, "text_frame"):
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ""
                        for run in paragraph.runs:
                            if run.text:
                                # 尝试多种编码方式处理文本
                                run_text = self._safe_decode_text(run.text)
                                if run_text:
                                    para_text += run_text
                        if para_text.strip():
                            text_parts.append(para_text.strip())
                except Exception as e:
                    print(f"提取text_frame失败: {str(e)}")
                    # 如果提取段落失败，尝试其他方法
                    pass
            
            # 方法2: 直接获取text属性（备用方法）
            if not text_parts and hasattr(shape, "text") and shape.text:
                text = self._safe_decode_text(shape.text)
                if text and text.strip():
                    text_parts.append(text.strip())
            
            # 方法3: 处理表格
            if hasattr(shape, "table"):
                try:
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = self._extract_text_from_shape(cell)
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            text_parts.append(" | ".join(row_texts))
                except Exception as e:
                    print(f"提取表格失败: {str(e)}")
                    pass
            
            # 方法4: 处理组合形状
            if hasattr(shape, "shapes"):
                try:
                    for sub_shape in shape.shapes:
                        sub_text = self._extract_text_from_shape(sub_shape)
                        if sub_text:
                            text_parts.append(sub_text)
                except Exception as e:
                    print(f"提取组合形状失败: {str(e)}")
                    pass
            
            # 方法5: 处理占位符文本
            if hasattr(shape, "placeholder_format"):
                try:
                    if hasattr(shape, "text") and shape.text:
                        placeholder_text = self._safe_decode_text(shape.text)
                        if placeholder_text and placeholder_text.strip():
                            text_parts.append(placeholder_text.strip())
                except Exception as e:
                    pass
                    
        except Exception as e:
            print(f"提取形状文本时出错: {str(e)}")
            # 如果提取失败，尝试基础方法
            try:
                if hasattr(shape, "text") and shape.text:
                    text = self._safe_decode_text(shape.text)
                    if text and text.strip():
                        text_parts.append(text.strip())
            except:
                pass
        
        return "\n".join(text_parts) if text_parts else ""
    
    def _safe_decode_text(self, text) -> str:
        """安全地解码文本，处理各种编码问题"""
        if not text:
            return ""
        
        try:
            # 如果已经是字符串，直接返回
            if isinstance(text, str):
                # 检查是否有编码问题
                try:
                    # 尝试编码为UTF-8，如果失败说明有编码问题
                    text.encode('utf-8')
                    return text
                except UnicodeEncodeError:
                    # 如果有编码问题，尝试修复
                    return text.encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')
            
            # 如果是字节，尝试解码
            elif isinstance(text, bytes):
                # 尝试多种编码
                encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1', 'cp1252']
                for encoding in encodings:
                    try:
                        return text.decode(encoding)
                    except UnicodeDecodeError:
                        continue
                # 如果都失败，使用错误处理
                return text.decode('utf-8', errors='ignore')
            
            # 其他类型，转换为字符串
            else:
                return str(text)
                
        except Exception as e:
            print(f"解码文本失败: {str(e)}")
            # 最后的备用方法
            try:
                return str(text).encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')
            except:
                return ""
    
    def _extract_text_fallback(self, prs: Presentation) -> str:
        """备用文本提取方法（当主要方法失败时使用）"""
        text_parts = []
        
        try:
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                
                # 尝试遍历所有可能的文本来源
                for shape in slide.shapes:
                    # 尝试所有可能的属性
                    for attr_name in ['text', 'text_frame', 'paragraphs']:
                        try:
                            if hasattr(shape, attr_name):
                                attr_value = getattr(shape, attr_name)
                                if attr_value:
                                    if attr_name == 'text':
                                        text = self._safe_decode_text(attr_value)
                                        if text and text.strip():
                                            slide_texts.append(text.strip())
                                    elif attr_name == 'text_frame':
                                        for para in attr_value.paragraphs:
                                            para_text = ""
                                            for run in para.runs:
                                                if run.text:
                                                    para_text += self._safe_decode_text(run.text)
                                            if para_text.strip():
                                                slide_texts.append(para_text.strip())
                        except Exception as e:
                            continue
                
                if slide_texts:
                    text_parts.append(f"\n--- 幻灯片 {slide_idx + 1} ---\n")
                    text_parts.extend(slide_texts)
        
        except Exception as e:
            print(f"备用提取方法也失败: {str(e)}")
        
        return "\n".join(text_parts)
    
    def _clean_text(self, text: str) -> str:
        """清理文本，移除控制字符和无效字符（增强版）"""
        if not text:
            return ""
        
        try:
            # 确保是字符串并正确编码
            if not isinstance(text, str):
                text = str(text)
            
            # 先进行安全解码
            text = self._safe_decode_text(text)
            
            # 移除控制字符（保留换行符和制表符）
            import re
            # 保留可打印字符、换行符、制表符、中文字符
            # 使用更宽松的规则，保留更多字符
            cleaned = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F-\x9F]', '', text)
            
            # 移除零宽字符
            cleaned = re.sub(r'[\u200B-\u200D\uFEFF]', '', cleaned)
            
            # 规范化空白字符
            cleaned = re.sub(r'[ \t]+', ' ', cleaned)  # 多个空格/制表符合并为一个空格
            cleaned = re.sub(r'\n\s*\n\s*\n+', '\n\n', cleaned)  # 多个换行符合并为两个
            
            # 移除BOM标记
            if cleaned.startswith('\ufeff'):
                cleaned = cleaned[1:]
            
            # 移除其他不可见字符，但保留中文、英文、数字和常用标点
            # 保留：中文、英文、数字、常用标点符号
            cleaned = re.sub(r'[^\u4e00-\u9fa5\w\s\.,;:!?()\[\]{}""''、。，；：！？（）【】《》\-\+\*/=]', '', cleaned)
            
            # 最终清理：移除多余的空白
            cleaned = cleaned.strip()
            
            return cleaned
        except Exception as e:
            print(f"清理文本失败: {str(e)}")
            # 如果清理失败，尝试最基本的清理
            try:
                return str(text).encode('utf-8', errors='ignore').decode('utf-8', errors='ignore').strip()
            except:
                return str(text) if text else ""

