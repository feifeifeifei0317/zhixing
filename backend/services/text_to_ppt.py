from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import re
import asyncio
from pathlib import Path
import tempfile


class TextToPPTService:
    """文本转PPT服务 - 精美版"""
    
    # 专业配色方案
    COLOR_SCHEME = {
        'primary': RGBColor(74, 144, 226),      # 主色：蓝色
        'secondary': RGBColor(52, 152, 219),    # 次色：浅蓝
        'accent': RGBColor(46, 204, 113),       # 强调色：绿色
        'text_dark': RGBColor(44, 62, 80),      # 深色文字
        'text_light': RGBColor(127, 140, 141),  # 浅色文字
        'background': RGBColor(247, 249, 250),  # 背景色
        'gradient_start': RGBColor(74, 144, 226),
        'gradient_end': RGBColor(52, 152, 219),
    }
    
    def __init__(self):
        self.temp_dir = Path("temp")
        self.temp_dir.mkdir(exist_ok=True)
    
    async def generate_ppt(self, text: str) -> str:
        """
        将文本转换为PPT
        
        Args:
            text: 输入的文本内容
            
        Returns:
            PPT文件路径
        """
        # 在后台线程中执行同步操作
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._generate_ppt_sync, text)
    
    def _generate_ppt_sync(self, text: str) -> str:
        """同步生成PPT - 精美版"""
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 解析文本内容
        sections = self._parse_text(text)
        
        # 生成精美的标题页
        self._create_title_slide(prs, sections[0]['title'] if sections else "课程演示")
        
        # 为每个章节创建精美的幻灯片
        for section in sections:
            # 创建章节内容幻灯片
            self._create_section_slides(prs, section)
        
        # 添加精美的结束页
        self._create_end_slide(prs)
        
        # 保存PPT
        output_path = self.temp_dir / f"presentation_{id(text)}.pptx"
        prs.save(str(output_path))
        
        return str(output_path)
    
    def _parse_text(self, text: str) -> list:
        """
        解析文本内容，提取章节和要点
        
        Returns:
            章节列表，每个章节包含title和points
        """
        sections = []
        lines = text.strip().split('\n')
        
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 检测章节标题（以"第X章"、"第X节"或数字开头）
            chapter_match = re.match(r'^(第[一二三四五六七八九十\d]+[章节]|第\d+[章节]|第\d+\.\d+[章节]|第\d+章|第\d+节|第\d+讲|第\d+部分)[：:]\s*(.+)', line)
            if chapter_match:
                if current_section:
                    sections.append(current_section)
                current_section = {
                    'title': chapter_match.group(2) or chapter_match.group(0),
                    'points': []
                }
                continue
            
            # 检测要点（以数字、字母或符号开头）
            point_match = re.match(r'^[\d\.\-\*\•\u2022]\s*(.+)', line)
            if point_match:
                if current_section:
                    current_section['points'].append(point_match.group(1))
                else:
                    # 如果没有当前章节，创建一个默认章节
                    current_section = {
                        'title': '课程内容',
                        'points': [point_match.group(1)]
                    }
            elif current_section:
                # 如果不是要点格式，但属于当前章节，也添加为要点
                if len(line) > 5:  # 过滤太短的行
                    current_section['points'].append(line)
        
        # 添加最后一个章节
        if current_section:
            sections.append(current_section)
        
        # 如果没有解析到章节，将整个文本作为一个章节
        if not sections:
            sections.append({
                'title': '课程内容',
                'points': [line.strip() for line in lines if line.strip()][:10]  # 最多10个要点
            })
        
        return sections
    
    def _create_title_slide(self, prs: Presentation, title_text: str):
        """创建精美的标题页"""
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)
        
        # 添加渐变背景矩形
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.COLOR_SCHEME['gradient_start']
        bg_shape.line.fill.background()
        
        # 添加装饰性圆形
        for i in range(3):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(7 + i * 0.5), Inches(0.5 + i * 0.3),
                Inches(2), Inches(2)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(
                255 - i * 20, 255 - i * 20, 255 - i * 20
            )
            circle.fill.transparency = 0.9
            circle.line.fill.background()
        
        # 添加主标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.LEFT
        title_para.space_after = Pt(20)
        
        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "AI智能生成 · 专业课程演示"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(230, 230, 230)
        subtitle_para.alignment = PP_ALIGN.LEFT
    
    def _create_section_slides(self, prs: Presentation, section: dict):
        """为章节创建精美的幻灯片"""
        section_title = section['title']
        points = section['points']
        
        # 每页最多显示6个要点，如果超过则分页
        max_points_per_slide = 6
        total_slides = (len(points) + max_points_per_slide - 1) // max_points_per_slide
        
        for slide_idx in range(total_slides):
            start_idx = slide_idx * max_points_per_slide
            end_idx = min(start_idx + max_points_per_slide, len(points))
            current_points = points[start_idx:end_idx]
            
            # 创建空白幻灯片
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # 添加背景色
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                prs.slide_width, prs.slide_height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.COLOR_SCHEME['background']
            bg_shape.line.fill.background()
            
            # 添加顶部装饰条
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                prs.slide_width, Inches(1.2)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = self.COLOR_SCHEME['primary']
            header_shape.line.fill.background()
            
            # 添加章节标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            display_title = section_title if slide_idx == 0 else f"{section_title} (续 {slide_idx + 1})"
            title_frame.text = display_title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.LEFT
            
            # 添加内容区域背景
            content_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(1.8),
                Inches(9), Inches(5.2)
            )
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
            content_bg.line.color.rgb = RGBColor(220, 220, 220)
            content_bg.line.width = Pt(1)
            
            # 添加要点列表
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8.5), Inches(4.8)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            for idx, point in enumerate(current_points):
                if idx == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                # 添加要点符号和文本
                p.text = f"  {point}"
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = self.COLOR_SCHEME['text_dark']
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.3
                
                # 添加装饰性图标（使用Unicode字符）
                run = p.runs[0]
                if idx == 0:
                    run.text = f"✓ {point}"
                    run.font.color.rgb = self.COLOR_SCHEME['accent']
                else:
                    run.text = f"• {point}"
    
    def _create_end_slide(self, prs: Presentation):
        """创建精美的结束页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 添加渐变背景
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.COLOR_SCHEME['gradient_end']
        bg_shape.line.fill.background()
        
        # 添加装饰性圆形
        for i in range(5):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1 + i * 1.5), Inches(1 + i * 0.2),
                Inches(1.5), Inches(1.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle.fill.transparency = 0.85
            circle.line.fill.background()
        
        # 添加结束文字
        end_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(6), Inches(2.5)
        )
        end_frame = end_box.text_frame
        end_frame.text = "谢谢观看"
        end_para = end_frame.paragraphs[0]
        end_para.font.size = Pt(56)
        end_para.font.bold = True
        end_para.font.color.rgb = RGBColor(255, 255, 255)
        end_para.alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(4.5), Inches(6), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "期待您的反馈与建议"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(240, 240, 240)
        subtitle_para.alignment = PP_ALIGN.CENTER







