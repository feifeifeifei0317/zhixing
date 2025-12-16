from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
import tempfile
import os
from pathlib import Path
import asyncio
import subprocess
import imageio_ffmpeg
import io


class PPTToVideoService:
    """PPT转视频服务"""
    
    def __init__(self):
        self.temp_dir = Path("temp")
        self.temp_dir.mkdir(exist_ok=True)
        self.slide_duration = 5  # 每张幻灯片显示5秒
    
    async def convert_to_video(
        self,
        ppt_path: str,
        ai_avatar_enabled: bool = True
    ) -> str:
        """
        将PPT转换为视频
        
        Args:
            ppt_path: PPT文件路径
            ai_avatar_enabled: 是否启用AI数字形象
            
        Returns:
            视频文件路径
        """
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None,
            self._convert_to_video_sync,
            ppt_path,
            ai_avatar_enabled
        )
    
    def _convert_to_video_sync(
        self,
        ppt_path: str,
        ai_avatar_enabled: bool
    ) -> str:
        """同步转换PPT为视频"""
        slide_images = []
        final_video = None
        try:
            import traceback
            print(f"开始转换PPT: {ppt_path}")
            
            # 加载PPT
            prs = Presentation(ppt_path)
            print(f"PPT加载成功，共 {len(prs.slides)} 张幻灯片")
            
            # 将每张幻灯片转换为图片
            slide_images = []
            for i, slide in enumerate(prs.slides):
                print(f"转换幻灯片 {i+1}/{len(prs.slides)}")
                img_path = self._slide_to_image(slide, i)
                slide_images.append(img_path)
            
            print(f"图片生成完成，共 {len(slide_images)} 张")
            
            # 使用 ImageSequenceClip 将图片列表合成为视频，每张保持固定时长
            print("开始创建视频片段...")
            final_video = ImageSequenceClip(
                slide_images,
                durations=[self.slide_duration] * len(slide_images)
            )
            print("视频片段创建成功")
            
            # 调整视频尺寸（尝试多种方法）
            try:
                # 方法1: 尝试使用 vfx.all.resize
                from moviepy.video.fx.all import resize
                final_video = final_video.fx(resize, height=720)
                print("视频尺寸调整完成")
            except Exception as e1:
                try:
                    # 方法2: 尝试直接使用 resize
                    final_video = final_video.resize(height=720)
                    print("视频尺寸调整完成（方法2）")
                except Exception as e2:
                        print(f"[WARN] 调整视频尺寸失败，使用原始尺寸: {str(e1)} / {str(e2)}")
                    # 如果调整尺寸失败，继续使用原始尺寸
            
            # 如果启用AI数字形象，添加虚拟人讲解（这里使用占位符）
            if ai_avatar_enabled:
                final_video = self._add_ai_avatar(final_video)
            
            # 输出视频路径
            output_path = self.temp_dir / f"video_{os.path.basename(ppt_path)}.mp4"
            print(f"开始导出视频到: {output_path}")
            
            # 设置FFmpeg路径（使用imageio_ffmpeg）
            ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
            os.environ['IMAGEIO_FFMPEG_EXE'] = ffmpeg_path
            
            # 导出视频（无音频版本，避免音频编解码器问题）
            final_video.write_videofile(
                str(output_path),
                fps=24,
                codec='libx264',
                audio_codec=None,  # 不使用音频编解码器
                remove_temp=True,
                ffmpeg_params=['-y']  # 覆盖输出文件
            )
            print(f"视频导出成功: {output_path}")
            
            # 清理临时图片文件
            for img_path in slide_images:
                if os.path.exists(img_path):
                    try:
                        os.unlink(img_path)
                    except:
                        pass
            
            return str(output_path)
            
        except Exception as e:
            import traceback
            error_trace = traceback.format_exc()
            print(f"转换视频失败: {str(e)}")
            print(f"错误堆栈: {error_trace}")
            
            # 清理资源
            if final_video:
                try:
                    final_video.close()
                except:
                    pass
            
            # 清理临时图片文件
            for img_path in slide_images:
                if os.path.exists(img_path):
                    try:
                        os.unlink(img_path)
                    except:
                        pass
            
            raise Exception(f"转换视频失败: {str(e)}")
    
    def _slide_to_image(self, slide, slide_index: int) -> str:
        """
        将幻灯片转换为图片（改进版，支持更好的文本提取）
        
        Args:
            slide: PPT幻灯片对象
            slide_index: 幻灯片索引
            
        Returns:
            图片文件路径
        """
        # 创建图片（1920x1080）
        width, height = 1920, 1080
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # 尝试加载字体（支持中文字体）
        try:
            # 尝试使用Windows系统字体
            import platform
            if platform.system() == 'Windows':
                try:
                    title_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 60)  # 微软雅黑
                    text_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 40)
                except:
                    try:
                        title_font = ImageFont.truetype("C:/Windows/Fonts/simsun.ttc", 60)  # 宋体
                        text_font = ImageFont.truetype("C:/Windows/Fonts/simsun.ttc", 40)
                    except:
                        title_font = ImageFont.load_default()
                        text_font = ImageFont.load_default()
            else:
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()
        except:
            title_font = ImageFont.load_default()
            text_font = ImageFont.load_default()
        
        # 安全提取文本的函数
        def safe_extract_text(shape):
            """安全地提取形状中的文本"""
            text_parts = []
            try:
                # 方法1: 从text_frame提取（最完整）
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ""
                        for run in paragraph.runs:
                            if run.text:
                                para_text += self._safe_decode_text(run.text)
                        if para_text.strip():
                            text_parts.append(para_text.strip())
                
                # 方法2: 直接获取text属性
                if not text_parts and hasattr(shape, "text") and shape.text:
                    text = self._safe_decode_text(shape.text)
                    if text and text.strip():
                        text_parts.append(text.strip())
            except Exception as e:
                # 如果提取失败，尝试基础方法
                try:
                    if hasattr(shape, "text") and shape.text:
                        text = self._safe_decode_text(shape.text)
                        if text:
                            text_parts.append(text)
                except:
                    pass
            return "\n".join(text_parts)
        
        # 先提取所有图片（在绘制之前）
        # 提取并绘制图片（改进版：保持原比例，优化布局）
        images_info = []  # 存储图片信息 [(image_obj, x, y, width, height, bbox)]
        image_rects = []  # 存储图片区域，用于文字布局避让 [(left, top, right, bottom)]
        
        # PPT默认尺寸（英寸）
        ppt_width_inches = 10
        ppt_height_inches = 7.5
        
        # 判断图片区域是否重叠（带安全间距）
        def is_overlap(x, y, w, h, rects, margin=16):
            """检查当前图片区域是否与已有图片重叠，带安全边距"""
            for rect in rects:
                left, top, right, bottom = rect
                if not (x + w + margin <= left or x >= right + margin or y + h + margin <= top or y >= bottom + margin):
                    return True
            return False

        # 为图片寻找不重叠的位置（优先目标区域，必要时平移/轻量缩放）
        def find_non_overlapping_position(px, py, w, h, rects, canvas_w, canvas_h, original_w, original_h):
            """尝试在画布内找到与已有图片不重叠的位置"""
            margin = 16
            scale = 1.0
            max_attempts = 12
            x, y = px, py

            def shrink():
                nonlocal w, h, scale, x, y
                scale *= 0.9
                w = max(1, int(original_w * scale))
                h = max(1, int(original_h * scale))
                x = min(max(0, x), max(0, canvas_w - w))
                y = min(max(0, y), max(0, canvas_h - h))

            attempt = 0
            while attempt < max_attempts:
                # 已经不重叠，返回
                if not is_overlap(x, y, w, h, rects, margin):
                    return x, y, w, h

                # 尝试向下移动到重叠图片的底部
                moved = False
                for left, top, right, bottom in rects:
                    if not (x + w + margin <= left or x >= right + margin or y + h + margin <= top or y >= bottom + margin):
                        # 先尝试向下挪
                        candidate_y = bottom + margin
                        if candidate_y + h <= canvas_h:
                            y = candidate_y
                            moved = True
                            break
                        # 向右挪
                        candidate_x = right + margin
                        if candidate_x + w <= canvas_w:
                            x = candidate_x
                            moved = True
                            break
                # 如果平移仍然重叠，则缩放
                if moved and not is_overlap(x, y, w, h, rects, margin):
                    return x, y, w, h
                if is_overlap(x, y, w, h, rects, margin):
                    shrink()
                attempt += 1

            # 兜底：保证在画布内返回当前位置
            x = min(max(0, x), max(0, canvas_w - w))
            y = min(max(0, y), max(0, canvas_h - h))
            return x, y, w, h

        # 提取所有图片
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    # 获取图片数据
                    image_stream = shape.image.blob
                    
                    # 获取形状的位置和尺寸（以EMU为单位）
                    left_emu = shape.left
                    top_emu = shape.top
                    width_emu = shape.width
                    height_emu = shape.height
                    
                    # 转换为英寸
                    left_inches = left_emu / 914400
                    top_inches = top_emu / 914400
                    width_inches = width_emu / 914400
                    height_inches = height_emu / 914400
                    
                    # 计算缩放比例（从PPT尺寸到我们的画布尺寸）
                    scale_x = width / (ppt_width_inches * 96)  # 假设96 DPI
                    scale_y = height / (ppt_height_inches * 96)
                    
                    # 转换为像素位置和尺寸（目标尺寸）
                    target_left = int(left_inches * 96 * scale_x)
                    target_top = int(top_inches * 96 * scale_y)
                    target_width = int(width_inches * 96 * scale_x)
                    target_height = int(height_inches * 96 * scale_y)
                    
                    # 加载原始图片
                    image_obj = Image.open(io.BytesIO(image_stream))
                    original_width, original_height = image_obj.size
                    
                    # 计算保持原比例的缩放尺寸
                    # 使用较小的缩放比例，确保图片完全适应目标区域
                    scale_w = target_width / original_width
                    scale_h = target_height / original_height
                    scale = min(scale_w, scale_h)  # 使用较小的比例保持原比例
                    
                    # 计算实际尺寸（保持原比例）
                    actual_width = int(original_width * scale)
                    actual_height = int(original_height * scale)
                    
                    # 初始放置：目标区域内居中
                    initial_left = target_left + (target_width - actual_width) // 2
                    initial_top = target_top + (target_height - actual_height) // 2
                    
                    # 确保不超出边界
                    initial_left = min(max(0, initial_left), max(0, width - actual_width))
                    initial_top = min(max(0, initial_top), max(0, height - actual_height))
                    
                    # 寻找不与其他图片重叠的位置（必要时轻量缩放）
                    adjusted_left, adjusted_top, adjusted_w, adjusted_h = find_non_overlapping_position(
                        initial_left,
                        initial_top,
                        actual_width,
                        actual_height,
                        image_rects,
                        width,
                        height,
                        original_width,
                        original_height
                    )
                    
                    # 确保尺寸有效
                    if adjusted_w > 0 and adjusted_h > 0:
                        # 调整图片大小（保持原比例）
                        image_obj = image_obj.resize((adjusted_w, adjusted_h), Image.Resampling.LANCZOS)
                        images_info.append((image_obj, adjusted_left, adjusted_top, adjusted_w, adjusted_h))
                        # 记录图片区域（用于文字避让和图片间避让）
                        image_rects.append((adjusted_left, adjusted_top, adjusted_left + adjusted_w, adjusted_top + adjusted_h))
                        print(f"提取图片: 位置({adjusted_left}, {adjusted_top}), 尺寸({adjusted_w}x{adjusted_h}), 原尺寸({original_width}x{original_height})")
                except Exception as e:
                    print(f"提取图片失败: {str(e)}")
                    import traceback
                    traceback.print_exc()
        
        # 先绘制所有图片（作为背景层）
        for image_obj, x, y, img_w, img_h in images_info:
            try:
                # 如果图片有透明通道，需要处理
                if image_obj.mode == 'RGBA':
                    # RGBA模式：使用alpha通道
                    img.paste(image_obj, (x, y), image_obj)
                elif image_obj.mode == 'P' and 'transparency' in image_obj.info:
                    # 调色板模式且有透明度信息
                    img.paste(image_obj, (x, y), image_obj.convert('RGBA').split()[3])
                else:
                    # 无透明通道，直接粘贴
                    img.paste(image_obj, (x, y))
            except Exception as e:
                print(f"绘制图片失败: {str(e)}")
                import traceback
                traceback.print_exc()
        
        # 绘制标题
        y_position = 100
        if slide.shapes.title:
            title_text = safe_extract_text(slide.shapes.title)
            if title_text:
                title_text = title_text.strip()
                # 计算文本位置（居中）
                try:
                    bbox = draw.textbbox((0, 0), title_text, font=title_font)
                    text_width = bbox[2] - bbox[0]
                    x_position = (width - text_width) // 2
                except:
                    x_position = width // 2
                
                draw.text(
                    (x_position, y_position),
                    title_text,
                    fill=(102, 126, 234),
                    font=title_font
                )
                y_position += 120
        
        
        # 检查点是否在图片区域内
        def is_point_in_image(x, y, image_rects):
            """检查点是否在任何图片区域内"""
            for rect in image_rects:
                left, top, right, bottom = rect
                if left <= x <= right and top <= y <= bottom:
                    return True
            return False
        
        def get_text_position(y_pos, text_height, image_rects, start_x=150, max_x=None):
            """获取文本位置，避开图片区域"""
            if max_x is None:
                max_x = width - 150
            
            # 检查当前位置是否与图片重叠
            text_bottom = y_pos + text_height
            for rect in image_rects:
                left, top, right, bottom = rect
                # 如果文本区域与图片区域重叠
                if not (text_bottom < top or y_pos > bottom):
                    # 文本与图片在垂直方向重叠，检查水平方向
                    if start_x < right and start_x + 50 > left:
                        # 有重叠，尝试调整位置
                        # 如果图片在左侧，文本向右移动
                        if right < width / 2:
                            return max(right + 20, start_x), y_pos
                        # 如果图片在右侧，文本保持左侧但限制宽度
                        else:
                            return start_x, min(right - 20, max_x)
            
            return start_x, y_pos
        
        # 绘制内容（改进的文本提取，避开图片区域）
        for shape in slide.shapes:
            if shape != slide.shapes.title:
                # 跳过图片形状（图片已经单独处理）
                if hasattr(shape, "image"):
                    continue
                    
                shape_text = safe_extract_text(shape)
                if shape_text:
                    lines = shape_text.split('\n')
                    for line in lines:
                        if line.strip():
                            # 处理项目符号
                            line = line.strip()
                            if not line.startswith(('•', '-', '*', '·')):
                                line = '• ' + line
                            
                            # 换行处理
                            if y_position > height - 100:
                                break
                            
                            # 计算文本位置，避开图片
                            text_x, text_y = get_text_position(y_position, 50, image_rects)
                            max_text_width = width - text_x - 50  # 留出右边距
                            
                            # 文本换行（如果太长）
                            try:
                                bbox = draw.textbbox((0, 0), line, font=text_font)
                                text_width = bbox[2] - bbox[0]
                                if text_width > max_text_width:
                                    # 需要换行
                                    words = line.split()
                                    current_line = ""
                                    current_y = text_y
                                    for word in words:
                                        test_line = current_line + " " + word if current_line else word
                                        test_bbox = draw.textbbox((0, 0), test_line, font=text_font)
                                        if test_bbox[2] - test_bbox[0] > max_text_width:
                                            if current_line:
                                                # 再次检查位置，避开图片
                                                final_x, final_y = get_text_position(current_y, 50, image_rects, text_x, width - 50)
                                                draw.text(
                                                    (final_x, final_y),
                                                    current_line,
                                                    fill=(50, 50, 50),
                                                    font=text_font
                                                )
                                                current_y += 50
                                            current_line = word
                                        else:
                                            current_line = test_line
                                    if current_line:
                                        final_x, final_y = get_text_position(current_y, 50, image_rects, text_x, width - 50)
                                        draw.text(
                                            (final_x, final_y),
                                            current_line,
                                            fill=(50, 50, 50),
                                            font=text_font
                                        )
                                        current_y += 50
                                    y_position = current_y
                                else:
                                    # 单行文本
                                    draw.text(
                                        (text_x, text_y),
                                        line,
                                        fill=(50, 50, 50),
                                        font=text_font
                                    )
                                    y_position += 50
                            except:
                                # 如果计算失败，直接绘制
                                draw.text(
                                    (text_x, text_y),
                                    line[:80],
                                    fill=(50, 50, 50),
                                    font=text_font
                                )
                                y_position += 50
        
        # 保存图片
        img_path = self.temp_dir / f"slide_{slide_index}.png"
        img.save(img_path)
        
        return str(img_path)
    
    def _safe_decode_text(self, text) -> str:
        """安全地解码文本，处理各种编码问题"""
        if not text:
            return ""
        
        try:
            # 如果已经是字符串，直接返回
            if isinstance(text, str):
                # 检查是否有编码问题
                try:
                    text.encode('utf-8')
                    return text
                except UnicodeEncodeError:
                    return text.encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')
            
            # 如果是字节，尝试解码
            elif isinstance(text, bytes):
                encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1', 'cp1252']
                for encoding in encodings:
                    try:
                        return text.decode(encoding)
                    except UnicodeDecodeError:
                        continue
                return text.decode('utf-8', errors='ignore')
            
            # 其他类型，转换为字符串
            else:
                return str(text)
                
        except Exception as e:
            try:
                return str(text).encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')
            except:
                return ""
    
    def _add_ai_avatar(self, video_clip):
        """
        添加AI数字形象到视频中
        
        注意：这是一个占位符实现。实际应用中需要集成真实的AI数字人服务，
        例如：腾讯云数字人、阿里云数字人、或者使用开源方案如SadTalker等。
        
        Args:
            video_clip: 视频片段
            
        Returns:
            添加了AI数字形象的视频片段
        """
        # TODO: 集成真实的AI数字人服务
        # 这里可以：
        # 1. 调用第三方API（如腾讯云、阿里云的数字人服务）
        # 2. 使用开源方案（如SadTalker）生成数字人视频
        # 3. 将数字人视频与PPT视频合成
        
        # 当前实现：返回原视频（占位符）
        # 实际应用中，这里应该：
        # - 提取PPT中的文本内容
        # - 调用TTS服务生成语音
        # - 调用数字人服务生成讲解视频
        # - 将数字人视频与PPT视频合成
        
        print("AI数字形象功能已启用（当前为占位符实现）")
        print("提示：需要集成真实的AI数字人服务API")
        
        return video_clip


